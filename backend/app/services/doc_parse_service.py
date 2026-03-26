import io
from app.services.ocr_service import OcrService


class DocParseService:
    def __init__(self):
        self.ocr = OcrService()

    def parse(self, content: bytes, file_type: str) -> str:
        """Return plain text extracted from document content."""
        ft = file_type.lower().lstrip(".")
        if ft == "pdf":
            return self._parse_pdf(content)
        elif ft == "docx":
            return self._parse_docx(content)
        elif ft == "pptx":
            return self._parse_pptx(content)
        elif ft in ("jpg", "jpeg", "png", "bmp"):
            return self.ocr.extract_text_from_image_bytes(content)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")

    def _parse_pdf(self, content: bytes) -> str:
        import fitz  # PyMuPDF
        doc = fitz.open(stream=content, filetype="pdf")
        pages = []
        for page in doc:
            text = page.get_text()
            if len(text.strip()) < 50:
                pix = page.get_pixmap()
                text = self.ocr.extract_text_from_image_bytes(pix.tobytes("png"))
            pages.append(text)
        return "\n".join(pages)

    def _parse_docx(self, content: bytes) -> str:
        from docx import Document as DocxDocument
        doc = DocxDocument(io.BytesIO(content))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

    def _parse_pptx(self, content: bytes) -> str:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(content))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        texts.append(para.text)
        return "\n".join(t for t in texts if t.strip())

    def _split_into_chunks(self, text: str, chunk_size: int = 500, overlap: int = 100) -> list[str]:
        """Split text into overlapping chunks."""
        if not text:
            return []
        chunks = []
        start = 0
        while start < len(text):
            end = min(start + chunk_size, len(text))
            chunks.append(text[start:end])
            if end == len(text):
                break
            start = end - overlap
        return chunks

    async def parse_task_documents(self, db, task_id: str) -> str:
        """Parse all documents for a task and return combined text."""
        from sqlalchemy import select
        from app.models.document import Document
        from app.services.file_service import FileService

        file_service = FileService()
        result = await db.execute(select(Document).where(Document.task_id == task_id))
        documents = result.scalars().all()
        texts = []
        for doc in documents:
            if doc.extracted_text:
                texts.append(doc.extracted_text)
            elif doc.storage_key:
                # Download from MinIO and parse
                content = await file_service.download_file(doc.storage_key)
                extracted = self.parse(content, doc.file_type)
                doc.extracted_text = extracted
                texts.append(extracted)
        # Persist extracted text so future calls don't re-parse
        await db.commit()
        return "\n\n".join(texts)
